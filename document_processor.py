# document_processor.py
from typing import Optional, Dict, Any
import fitz
import docx
import pptx
from pptx import Presentation
import io
import logging
from abc import ABC, abstractmethod

class DocumentProcessor(ABC):
    @abstractmethod
    def extract_text(self, file_obj: io.BytesIO) -> str:
        pass

class PDFProcessor(DocumentProcessor):
    def extract_text(self, file_obj: io.BytesIO) -> str:
        text = ""
        with fitz.open(stream=file_obj, filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text

class DocxProcessor(DocumentProcessor):
    def extract_text(self, file_obj: io.BytesIO) -> str:
        doc = docx.Document(file_obj)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

class TxtProcessor(DocumentProcessor):
    def extract_text(self, file_obj: io.BytesIO) -> str:
        return file_obj.read().decode('utf-8')

class PPTXProcessor(DocumentProcessor):
    def extract_text(self, file_obj: io.BytesIO) -> str:
        presentation = Presentation(file_obj)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

class DocumentFactory:
    _processors = {
        "application/pdf": PDFProcessor(),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocxProcessor(),
        "text/plain": TxtProcessor(),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": PPTXProcessor()
    }

    @classmethod
    def get_processor(cls, mime_type: str) -> Optional[DocumentProcessor]:
        return cls._processors.get(mime_type.lower())